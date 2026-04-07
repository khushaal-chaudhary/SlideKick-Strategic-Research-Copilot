"""
Generator Node - Creates deliverables in the appropriate format.

This node:
1. Takes the analyzed insights
2. Decides on the best output format (or uses the planned one)
3. Generates content in that format
4. For slides, creates Google Slides via API
"""

import json
import logging
import subprocess
from pathlib import Path
from typing import Any

from copilot.agent.state import OutputFormat, QueryType, ResearchState
from copilot.llm import get_llm

logger = logging.getLogger(__name__)


CHAT_RESPONSE_PROMPT = """You are a strategic research analyst providing findings.

## Original Query
{query}

## Key Insights
{insights}

## Full Synthesis
{synthesis}

## Quality Context
- Confidence Score: {quality_score:.0%}
- Entities Analyzed: {entity_count}

## Task
Provide a clear, professional response that:
1. Directly answers the query
2. Highlights the most important insights
3. Notes any limitations or gaps in the analysis
4. Is conversational but substantive

Keep the response focused and actionable. For strategic queries, include recommendations.
"""


SLIDES_CONTENT_PROMPT = """You are creating a structured slide deck for an executive presentation.

## Original Query
{query}

## Query Type
{query_type}

## Key Insights
{insights}

## Full Synthesis
{synthesis}

## Your Job
Pick a color theme, then build 5-8 slides using the typed layouts below. You decide WHAT content goes on each slide and WHICH layout best fits that content. The visual design is handled automatically — you just fill content slots.

## Theme Choices (pick ONE)
- "slate" — Dark navy + white. Best for corporate/executive presentations.
- "growth" — Deep green + gold. Best for financial/investment content.
- "clarity" — White + blue. Best for product/technology topics.
- "bold" — Black + electric accent. Best for startup/pitch energy.
- "warm" — Off-white + terracotta. Best for strategy/consulting.

## Available Layouts

### "title" — Opening slide
content: {{ "headline": "string", "subtitle": "string" }}

### "stat-trio" — Three key metrics side by side
content: {{ "title": "string", "stats": [ {{ "value": "$52B", "label": "Revenue", "delta": "+18% YoY" }}, ... ] }}
Rules: Exactly 3 stats. "value" should be short (number + unit). "delta" is a change indicator.

### "bullets" — Narrative content with rich bullet points
content: {{ "title": "string", "category": "strategy|financial_metric|risk|opportunity|technology|competition|growth|leadership", "bullets": [ {{ "heading": "2-5 word key point", "detail": "15-25 word explanation with data" }}, ... ] }}
Rules: 3-4 bullets max. Each bullet has a heading AND a detail.

### "two-column" — Side-by-side comparison
content: {{ "title": "string", "left": {{ "label": "string", "points": ["string", ...] }}, "right": {{ "label": "string", "points": ["string", ...] }} }}
Rules: 2-4 points per side. Use for comparisons, pros/cons, before/after.

### "quote-callout" — Highlight a key finding
content: {{ "title": "string", "category": "strategy|financial_metric|risk|opportunity|technology|competition|growth|leadership", "quote": "The key insight or finding (1-2 sentences)", "attribution": "Based on...", "supporting": "Additional context line" }}

### "timeline" — Chronological events
content: {{ "title": "string", "events": [ {{ "date": "2020", "label": "What happened (8-15 words)" }}, ... ] }}
Rules: 3-6 events. Dates can be years, quarters, or months.

### "closing" — Recommendations and next steps
content: {{ "title": "string", "actions": [ {{ "heading": "Immediate/Short-term/Strategic", "detail": "Specific actionable recommendation (15-25 words)" }}, ... ] }}
Rules: 2-4 actions.

## Slide Deck Guidelines
1. ALWAYS start with a "title" slide
2. ALWAYS end with a "closing" slide
3. Use "stat-trio" when you have quantitative metrics — executives love numbers
4. Use "quote-callout" for your single strongest insight
5. Use "two-column" for any comparative analysis
6. Use "timeline" if the query involves evolution or history
7. Use "bullets" for everything else — it's your workhorse
8. Each slide should have speaker_notes (2-3 sentences for the presenter)
9. For financial queries, prefer "growth" theme and lead with "stat-trio"
10. For strategic queries, prefer "slate" theme
11. For comparative queries, always include at least one "two-column"

## Response Format (JSON)
{{
    "theme": "slate",
    "title": "Presentation Title",
    "subtitle": "Subtitle",
    "slides": [
        {{
            "layout": "title",
            "content": {{ "headline": "...", "subtitle": "..." }},
            "speaker_notes": "..."
        }},
        {{
            "layout": "stat-trio",
            "content": {{ "title": "...", "stats": [...] }},
            "speaker_notes": "..."
        }}
    ]
}}

Respond with valid JSON only, no markdown formatting.
"""


BULLET_SUMMARY_PROMPT = """You are creating a structured summary.

## Original Query
{query}

## Key Insights
{insights}

## Synthesis
{synthesis}

## Task
Create a structured bullet summary with clear sections:

# [Topic/Query Summary]

## Key Findings
• Finding 1
• Finding 2
• Finding 3

## Details
### [Subtopic 1]
• Detail 1
• Detail 2

### [Subtopic 2]
• Detail 1
• Detail 2

## Gaps/Limitations
• What we couldn't find

## Recommendations (if applicable)
• Recommendation 1
• Recommendation 2

Keep it scannable and professional.
"""


def _format_insights_for_generation(insights: list[dict]) -> str:
    """Format insights for generation prompts."""
    if not insights:
        return "No specific insights were extracted."
    
    formatted = []
    for insight in insights:
        formatted.append(
            f"**{insight.get('title', 'Insight')}** ({insight.get('category', 'general')})\n"
            f"{insight.get('description', '')}\n"
            f"Evidence: {', '.join(insight.get('supporting_evidence', [])[:3])}\n"
            f"Confidence: {insight.get('confidence', 0):.0%}"
        )
    return "\n\n".join(formatted)


def _generate_chat_response(state: ResearchState) -> str:
    """Generate a conversational response."""
    llm = get_llm(temperature=0.3)
    
    prompt = CHAT_RESPONSE_PROMPT.format(
        query=state["original_query"],
        insights=_format_insights_for_generation(state.get("insights", [])),
        synthesis=state.get("synthesis", "No synthesis available."),
        quality_score=state.get("quality_score", 0.5),
        entity_count=len(state.get("entities_found", [])),
    )
    
    response = llm.invoke(prompt)
    return response.content


def _generate_bullet_summary(state: ResearchState) -> str:
    """Generate a structured bullet summary."""
    llm = get_llm(temperature=0.3)
    
    prompt = BULLET_SUMMARY_PROMPT.format(
        query=state["original_query"],
        insights=_format_insights_for_generation(state.get("insights", [])),
        synthesis=state.get("synthesis", ""),
    )
    
    response = llm.invoke(prompt)
    return response.content


def _generate_slides_content(state: ResearchState) -> dict[str, Any]:
    """Generate slide deck structure using typed layouts."""
    llm = get_llm(temperature=0.3)

    prompt = SLIDES_CONTENT_PROMPT.format(
        query=state["original_query"],
        query_type=state.get("query_type", "strategic"),
        insights=_format_insights_for_generation(state.get("insights", [])),
        synthesis=state.get("synthesis", ""),
    )

    response = llm.invoke(prompt)

    # Parse JSON response
    cleaned = response.content.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.split("```")[1]
        if cleaned.startswith("json"):
            cleaned = cleaned[4:]
    cleaned = cleaned.strip()

    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        # Try regex extraction
        import re
        match = re.search(r"\{.*\}", response.content, re.DOTALL)
        if match:
            try:
                data = json.loads(match.group())
            except json.JSONDecodeError:
                data = None
        else:
            data = None

    if not data or "slides" not in data:
        # Fallback structure using typed layouts
        synthesis = state.get("synthesis", "Analysis results are available.")
        data = {
            "theme": "slate",
            "title": "Research Findings",
            "subtitle": state["original_query"][:60],
            "slides": [
                {
                    "layout": "title",
                    "content": {"headline": "Research Findings", "subtitle": state["original_query"][:60]},
                    "speaker_notes": "Opening slide.",
                },
                {
                    "layout": "bullets",
                    "content": {
                        "title": "Summary",
                        "category": "strategy",
                        "bullets": [{"heading": "Key Finding", "detail": synthesis[:200]}],
                    },
                    "speaker_notes": "Main findings from the analysis.",
                },
                {
                    "layout": "closing",
                    "content": {
                        "title": "Next Steps",
                        "actions": [{"heading": "Review", "detail": "Review findings and determine follow-up research areas."}],
                    },
                    "speaker_notes": "Closing remarks.",
                },
            ],
        }

    # Ensure theme is set
    if "theme" not in data:
        data["theme"] = "slate"

    return data


def _convert_to_mcp_batch_requests(slides_content: dict, presentation_id: str) -> list[dict]:
    """
    Convert generator's slide format to Google Slides API batch update requests.

    Compatible with @bohachu/google-slides-mcp batch_update_presentation tool.
    """
    requests = []

    for idx, slide in enumerate(slides_content.get("slides", [])):
        slide_type = slide.get("type", "content")
        slide_object_id = f"slide_{idx}"

        # Determine layout based on slide type
        if slide_type == "title":
            layout = "TITLE"
        elif slide_type == "section":
            layout = "SECTION_HEADER"
        else:
            layout = "TITLE_AND_BODY"

        # Create slide request
        requests.append({
            "createSlide": {
                "objectId": slide_object_id,
                "insertionIndex": idx,
                "slideLayoutReference": {
                    "predefinedLayout": layout
                }
            }
        })

        # Add title text
        title_text = slide.get("title", "")
        if title_text:
            title_box_id = f"title_{idx}"
            requests.append({
                "createShape": {
                    "objectId": title_box_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": slide_object_id,
                        "size": {"width": {"magnitude": 600, "unit": "PT"}, "height": {"magnitude": 50, "unit": "PT"}},
                        "transform": {"scaleX": 1, "scaleY": 1, "translateX": 50, "translateY": 30, "unit": "PT"}
                    }
                }
            })
            requests.append({
                "insertText": {
                    "objectId": title_box_id,
                    "text": title_text,
                    "insertionIndex": 0
                }
            })

        # Add subtitle for title slides
        if slide_type == "title" and slide.get("subtitle"):
            subtitle_box_id = f"subtitle_{idx}"
            requests.append({
                "createShape": {
                    "objectId": subtitle_box_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": slide_object_id,
                        "size": {"width": {"magnitude": 500, "unit": "PT"}, "height": {"magnitude": 30, "unit": "PT"}},
                        "transform": {"scaleX": 1, "scaleY": 1, "translateX": 100, "translateY": 90, "unit": "PT"}
                    }
                }
            })
            requests.append({
                "insertText": {
                    "objectId": subtitle_box_id,
                    "text": slide.get("subtitle", ""),
                    "insertionIndex": 0
                }
            })

        # Add bullet points for content slides
        bullets = slide.get("bullets", [])
        if bullets and slide_type != "title":
            body_box_id = f"body_{idx}"
            bullet_text = "\n".join(f"• {b}" for b in bullets)
            requests.append({
                "createShape": {
                    "objectId": body_box_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": slide_object_id,
                        "size": {"width": {"magnitude": 600, "unit": "PT"}, "height": {"magnitude": 300, "unit": "PT"}},
                        "transform": {"scaleX": 1, "scaleY": 1, "translateX": 50, "translateY": 100, "unit": "PT"}
                    }
                }
            })
            requests.append({
                "insertText": {
                    "objectId": body_box_id,
                    "text": bullet_text,
                    "insertionIndex": 0
                }
            })

    return requests


def _get_service_account_credentials():
    """
    Load Google Service Account credentials.

    Looks for credentials in order:
    1. GOOGLE_SERVICE_ACCOUNT_JSON env var (JSON content directly)
    2. GOOGLE_APPLICATION_CREDENTIALS env var (path to JSON file)
    3. Local development path: packages/google-slides-mcp/keys/google_service_account_key.json

    Returns None if no credentials found (Google Slides is optional).
    """
    import json
    import os
    import tempfile

    from google.oauth2 import service_account

    SCOPES = [
        "https://www.googleapis.com/auth/presentations",
        "https://www.googleapis.com/auth/drive.file",
    ]

    # Option 1: JSON content directly in environment variable (best for HF Spaces)
    json_content = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON")
    if json_content:
        try:
            creds_dict = json.loads(json_content)
            return service_account.Credentials.from_service_account_info(
                creds_dict,
                scopes=SCOPES,
            )
        except json.JSONDecodeError as e:
            logger.warning("Failed to parse GOOGLE_SERVICE_ACCOUNT_JSON: %s", e)

    # Option 2: Path to credentials file via environment variable
    env_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
    if env_path and Path(env_path).exists():
        try:
            return service_account.Credentials.from_service_account_file(
                env_path,
                scopes=SCOPES,
            )
        except Exception as e:
            logger.warning("Failed to load credentials from %s: %s", env_path, e)

    # Option 3: Local development path (relative to project root)
    # From: packages/agent/src/copilot/agent/nodes/generator.py
    # To:   packages/google-slides-mcp/keys/google_service_account_key.json
    local_paths = [
        # Development: relative to generator.py
        Path(__file__).parent.parent.parent.parent.parent.parent
        / "google-slides-mcp"
        / "keys"
        / "google_service_account_key.json",
        # Docker/HF Spaces: might be in app root
        Path("/home/user/app/keys/google_service_account_key.json"),
        # Current working directory
        Path("keys/google_service_account_key.json"),
    ]

    for local_key_path in local_paths:
        if local_key_path.exists():
            try:
                return service_account.Credentials.from_service_account_file(
                    str(local_key_path),
                    scopes=SCOPES,
                )
            except Exception as e:
                logger.warning("Failed to load credentials from %s: %s", local_key_path, e)

    # No credentials found - this is OK, Google Slides is optional
    logger.debug("Google Service Account not configured - Slides generation disabled")
    return None


def _get_share_email() -> str | None:
    """
    Get the email address to share presentations with.

    Checks:
    1. Environment variable GOOGLE_SLIDES_SHARE_EMAIL
    2. Could be extended to check a config file
    """
    import os

    return os.environ.get("GOOGLE_SLIDES_SHARE_EMAIL")


def _share_presentation(presentation_id: str, email: str, creds) -> tuple[bool, str | None]:
    """
    Share a Google Slides presentation with a specific email.

    Args:
        presentation_id: The ID of the presentation to share
        email: The email address to share with
        creds: Google service account credentials

    Returns:
        Tuple of (success, error_message)
    """
    try:
        from googleapiclient.discovery import build

        drive_service = build("drive", "v3", credentials=creds)

        # Create permission to share with user
        permission = {
            "type": "user",
            "role": "writer",  # Can edit the presentation
            "emailAddress": email,
        }

        drive_service.permissions().create(
            fileId=presentation_id,
            body=permission,
            sendNotificationEmail=True,  # Send email notification
            emailMessage="A new research presentation has been shared with you from Strategic Research Copilot.",
        ).execute()

        logger.info("   Shared presentation with: %s", email)
        return True, None

    except Exception as e:
        logger.warning("   Failed to share presentation: %s", e)
        return False, str(e)


def _create_google_slides(
    slides_content: dict, share_email: str | None = None
) -> dict[str, Any]:
    """
    Create Google Slides presentation using Service Account authentication.

    Uses @bohachu/google-slides-mcp compatible approach with service account.

    Args:
        slides_content: The slide deck structure to create
        share_email: Optional email to share the presentation with (overrides config)

    Returns:
        Dict with keys:
        - url: Presentation URL (or None if failed)
        - error: Error message (or None if success)
        - presentation_id: The presentation ID
        - shared: Whether the presentation was shared
        - share_email_needed: True if email is needed for sharing
    """
    result = {
        "url": None,
        "error": None,
        "presentation_id": None,
        "shared": False,
        "share_email_needed": False,
    }

    try:
        from googleapiclient.discovery import build

        # Get service account credentials
        creds = _get_service_account_credentials()
        if not creds:
            result["error"] = (
                "Google Slides disabled - no service account configured. "
                "Set GOOGLE_SERVICE_ACCOUNT_JSON env var with JSON content, "
                "or GOOGLE_APPLICATION_CREDENTIALS with path to key file."
            )
            logger.info("   Skipping Google Slides - credentials not configured")
            return result

        # Build Slides API service
        slides_service = build("slides", "v1", credentials=creds)

        # Step 1: Create empty presentation
        presentation = (
            slides_service.presentations()
            .create(body={"title": slides_content.get("title", "Research Presentation")})
            .execute()
        )

        presentation_id = presentation.get("presentationId")
        url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"

        result["presentation_id"] = presentation_id
        result["url"] = url

        logger.info("   Created presentation: %s", presentation_id)

        # Step 2: Build batch update requests for slides content
        batch_requests = _convert_to_mcp_batch_requests(slides_content, presentation_id)

        if batch_requests:
            # Execute batch update to add all slides and content
            slides_service.presentations().batchUpdate(
                presentationId=presentation_id, body={"requests": batch_requests}
            ).execute()
            logger.info("   Added %d slides with content", len(slides_content.get("slides", [])))

        # Step 3: Share the presentation
        # Priority: provided email > configured email
        email_to_share = share_email or _get_share_email()

        if email_to_share:
            success, share_error = _share_presentation(presentation_id, email_to_share, creds)
            result["shared"] = success
            if not success:
                logger.warning("   Could not share: %s", share_error)
        else:
            # No email configured - flag that we need one
            result["share_email_needed"] = True
            logger.info("   No share email configured - user will need to provide one")

        return result

    except ImportError:
        result["error"] = "Google API not installed. Run: pip install google-api-python-client google-auth"
        return result
    except Exception as e:
        logger.exception("Failed to create Google Slides")
        result["error"] = f"Failed to create slides: {str(e)}"
        return result


def _create_pptx_presentation(slides_content: dict, session_id: str = None) -> dict[str, Any]:
    """
    Create a PowerPoint presentation using python-pptx as a fallback.

    This is used when Google Slides API is not available.
    Creates a basic but clean presentation that users can enhance.

    Args:
        slides_content: The slide deck structure
        session_id: Optional session ID for unique filename

    Returns:
        Dict with keys:
        - file_path: Path to the generated .pptx file
        - filename: Just the filename
        - error: Error message if failed
        - design_tips: Tips for enhancing the presentation
    """
    import os
    import tempfile
    import uuid

    result = {
        "file_path": None,
        "filename": None,
        "error": None,
        "design_tips": None,
    }

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Color scheme
        TITLE_COLOR = RGBColor(26, 54, 93)  # Dark blue #1a365d
        BODY_COLOR = RGBColor(45, 55, 72)   # Dark gray #2d3748
        ACCENT_COLOR = RGBColor(49, 130, 206)  # Blue accent #3182ce

        # Collect design suggestions for output
        all_design_tips = []

        for idx, slide_data in enumerate(slides_content.get("slides", [])):
            slide_type = slide_data.get("type", "content")

            # Use blank layout for full control
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Get design suggestions if available
            design_suggestions = slide_data.get("design_suggestions", {})
            if design_suggestions:
                tip_text = f"**Slide {idx + 1}: {slide_data.get('title', 'Untitled')}**\n"
                if design_suggestions.get("visual_type"):
                    tip_text += f"  • Visual: {design_suggestions.get('visual_type')}\n"
                if design_suggestions.get("visual_description"):
                    tip_text += f"  • Description: {design_suggestions.get('visual_description')}\n"
                if design_suggestions.get("layout_tip"):
                    tip_text += f"  • Layout: {design_suggestions.get('layout_tip')}\n"
                if design_suggestions.get("icons_to_use"):
                    tip_text += f"  • Icons: {design_suggestions.get('icons_to_use')}\n"
                all_design_tips.append(tip_text)

            if slide_type == "title":
                # Title slide - centered title and subtitle
                title_text = slide_data.get("title", "Presentation")
                subtitle_text = slide_data.get("subtitle", "")

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = title_text
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = TITLE_COLOR
                title_para.alignment = PP_ALIGN.CENTER

                # Subtitle
                if subtitle_text:
                    subtitle_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
                    )
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_para = subtitle_frame.paragraphs[0]
                    subtitle_para.text = subtitle_text
                    subtitle_para.font.size = Pt(24)
                    subtitle_para.font.color.rgb = BODY_COLOR
                    subtitle_para.alignment = PP_ALIGN.CENTER

            else:
                # Content slides (executive_summary, content, recommendations)
                title_text = slide_data.get("title", "")
                bullets = slide_data.get("bullets", [])
                speaker_notes = slide_data.get("speaker_notes", "")

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.4), Inches(12.333), Inches(1)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = title_text
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                title_para.font.color.rgb = TITLE_COLOR

                # Bullets
                if bullets:
                    body_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
                    )
                    body_frame = body_box.text_frame
                    body_frame.word_wrap = True

                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            para = body_frame.paragraphs[0]
                        else:
                            para = body_frame.add_paragraph()
                        para.text = f"• {bullet}"
                        para.font.size = Pt(18)  # Slightly smaller for longer content
                        para.font.color.rgb = BODY_COLOR
                        para.space_after = Pt(16)  # More spacing between bullets

                # Add speaker notes
                if speaker_notes:
                    notes_slide = slide.notes_slide
                    notes_frame = notes_slide.notes_text_frame
                    notes_frame.text = speaker_notes

                    # Add design suggestions to notes if available
                    if design_suggestions:
                        notes_frame.text += "\n\n--- DESIGN SUGGESTIONS ---\n"
                        if design_suggestions.get("visual_type"):
                            notes_frame.text += f"Visual: {design_suggestions.get('visual_type')}\n"
                        if design_suggestions.get("visual_description"):
                            notes_frame.text += f"Description: {design_suggestions.get('visual_description')}\n"
                        if design_suggestions.get("layout_tip"):
                            notes_frame.text += f"Layout: {design_suggestions.get('layout_tip')}\n"
                        if design_suggestions.get("color_emphasis"):
                            notes_frame.text += f"Colors: {design_suggestions.get('color_emphasis')}\n"
                        if design_suggestions.get("icons_to_use"):
                            notes_frame.text += f"Icons: {design_suggestions.get('icons_to_use')}\n"

        # Save to temp directory
        temp_dir = os.path.join(tempfile.gettempdir(), "slidekick_presentations")
        os.makedirs(temp_dir, exist_ok=True)

        # Generate unique filename
        file_id = session_id or str(uuid.uuid4())[:8]
        safe_title = "".join(c for c in slides_content.get("title", "presentation")[:30] if c.isalnum() or c in " -_").strip()
        safe_title = safe_title.replace(" ", "_") or "presentation"
        filename = f"{safe_title}_{file_id}.pptx"
        file_path = os.path.join(temp_dir, filename)

        # Save presentation
        prs.save(file_path)

        result["file_path"] = file_path
        result["filename"] = filename
        result["design_tips"] = "\n".join(all_design_tips) if all_design_tips else None

        logger.info("   Created PPTX presentation: %s", filename)
        return result

    except ImportError as e:
        # Log the actual import error for debugging
        logger.error("ImportError in PPTX generation: %s", str(e))
        # Check if it's specifically python-pptx or a dependency
        error_msg = str(e).lower()
        if "pptx" in error_msg:
            result["error"] = "python-pptx not installed. Run: pip install python-pptx"
        elif "lxml" in error_msg:
            result["error"] = f"Missing dependency for python-pptx: {e}. lxml may need to be installed."
        else:
            result["error"] = f"Import error during presentation creation: {e}"
        return result
    except Exception as e:
        logger.exception("Failed to create PPTX presentation")
        result["error"] = f"Failed to create presentation: {str(e)}"
        return result


def generator_node(state: ResearchState) -> dict[str, Any]:
    """
    Generate the final deliverable.
    
    This node:
    1. Determines the output format
    2. Generates content appropriate for that format
    3. For slides, attempts to create via MCP
    4. Returns the generated content
    
    Returns:
        State updates with generated content
    """
    output_format = state.get("output_format", OutputFormat.CHAT.value)
    query_type = state.get("query_type", QueryType.UNKNOWN.value)
    
    logger.info("🎨 Generator: Creating %s output...", output_format)

    output_content = ""
    output_url = None
    slides_content = None

    if output_format == OutputFormat.CHAT.value:
        output_content = _generate_chat_response(state)
        
    elif output_format == OutputFormat.BULLET_SUMMARY.value:
        output_content = _generate_bullet_summary(state)
        
    elif output_format == OutputFormat.SLIDES.value:
        slides_content = _generate_slides_content(state)

        # Check if user provided an email for sharing (from previous interaction)
        user_share_email = state.get("user_share_email")

        # Try to create actual Google Slides first
        google_slides_url = None
        try:
            result = _create_google_slides(slides_content, share_email=user_share_email)
            google_slides_url = result.get("url")
        except Exception as e:
            logger.info("   Google Slides unavailable: %s", e)

        if google_slides_url:
            output_url = google_slides_url
            output_content = f"✅ Created presentation: {google_slides_url}\n\n"
            output_content += f"**{slides_content.get('title', 'Presentation')}**\n\n"
            output_content += "Slides:\n"
            for i, slide in enumerate(slides_content.get("slides", []), 1):
                title = slide.get("content", {}).get("title", slide.get("content", {}).get("headline", "Untitled"))
                output_content += f"  {i}. {title}\n"

            if result.get("shared"):
                share_email = user_share_email or _get_share_email()
                output_content += f"\n✅ Shared with: {share_email}\n"
        else:
            # Default: send slides JSON to frontend for Reveal.js rendering
            logger.info("   Using browser presentation (Reveal.js)")
            output_content = f"**{slides_content.get('title', 'Presentation')}**\n"
            output_content += f"_{slides_content.get('subtitle', '')}_\n\n"
            output_content += f"{len(slides_content.get('slides', []))} slides generated. "
            output_content += "View the interactive presentation above."
    
    elif output_format == OutputFormat.DOCUMENT.value:
        # For now, generate as extended chat response
        output_content = _generate_chat_response(state)
        output_content = "# Research Report\n\n" + output_content
    
    logger.info("   Generated %d characters of content", len(output_content))

    result = {
        "output_content": output_content,
        "output_url": output_url,
    }

    # Pass slides data to frontend for Reveal.js rendering
    if output_format == OutputFormat.SLIDES.value and slides_content:
        result["slides_content"] = slides_content

    return result